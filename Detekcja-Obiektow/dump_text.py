"""Dump tekstowy zawartości slajdów - weryfikacja treści."""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

p = Presentation("D:/CLCO/grupa2_detekcja/GRUPA2_Detekcja_Obiektow_Postepy.pptx")

for idx, slide in enumerate(p.slides, 1):
    print(f"\n{'='*70}\nSLAJD {idx}\n{'='*70}")
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                first_line = txt.split('\n')[0][:80]
                full = txt[:200].replace('\n', ' | ')
                print(f"  • {full}{'...' if len(txt) > 200 else ''}")
