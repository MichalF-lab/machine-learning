"""QA prezentacji - sprawdzenie zawartości i wymiarów elementów."""
from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400

def emu_to_in(emu):
    return emu / EMU_PER_INCH

p = Presentation("D:/CLCO/grupa2_detekcja/GRUPA2_Detekcja_Obiektow_Postepy.pptx")

print(f"Prezentacja: {len(p.slides)} slajdów")
print(f"Wymiary slajdu: {emu_to_in(p.slide_width):.2f} x {emu_to_in(p.slide_height):.2f} cali")
print("=" * 80)

SLIDE_W = emu_to_in(p.slide_width)
SLIDE_H = emu_to_in(p.slide_height)

issues = []

for idx, slide in enumerate(p.slides, 1):
    print(f"\n--- SLAJD {idx} ---")
    shape_count = len(slide.shapes)
    print(f"  Liczba kształtów: {shape_count}")

    for sh in slide.shapes:
        try:
            x, y = emu_to_in(sh.left), emu_to_in(sh.top)
            w, h = emu_to_in(sh.width), emu_to_in(sh.height)

            # Sprawdź czy element wystaje poza slajd
            if x < -0.05 or y < -0.05:
                issues.append(f"Slajd {idx}: kształt poza lewym/górnym brzegiem ({x:.2f}, {y:.2f}) - {sh.shape_type}")
            if x + w > SLIDE_W + 0.05:
                issues.append(f"Slajd {idx}: kształt wystaje w prawo ({x+w:.2f} > {SLIDE_W}) - {sh.shape_type}")
            if y + h > SLIDE_H + 0.05:
                issues.append(f"Slajd {idx}: kształt wystaje w dół ({y+h:.2f} > {SLIDE_H}) - {sh.shape_type}")

            # Sprawdź tekst
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if txt and len(txt) > 60:
                    # Sprawdź czy tekst nie jest za duży na pole
                    if h < 0.3 and len(txt) > 80:
                        issues.append(f"Slajd {idx}: małe pole z dużym tekstem ({h:.2f}\" h, {len(txt)} chars): '{txt[:50]}...'")
                # Sprawdź placeholder/lorem
                low = txt.lower()
                if "lorem" in low or "ipsum" in low or "todo" in low or "[insert" in low:
                    issues.append(f"Slajd {idx}: placeholder w tekście: '{txt[:50]}'")
        except (AttributeError, TypeError):
            pass

print("\n" + "=" * 80)
print(f"ZNALEZIONE PROBLEMY: {len(issues)}")
for i in issues:
    print(f"  ⚠ {i}")
if not issues:
    print("  ✓ brak problemów geometrycznych / tekstowych")
